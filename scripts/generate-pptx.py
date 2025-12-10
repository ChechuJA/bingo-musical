#!/usr/bin/env python3
"""
Generador de PowerPoint para Cartones de Bingo Musical
Crea presentaciones con 3 cartones por diapositiva
Requiere: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def create_bingo_pptx(cards, categoria, theme_colors, output_file):
    """
    Crea presentación PowerPoint con cartones de bingo
    Layout dinámico según número de canciones:
    - 6-8 canciones: 3 cartones por slide
    - 10-12 canciones: 2 cartones por slide
    - 15-20 canciones: 1 cartón por slide
    
    Args:
        cards: Lista de cartones (cada cartón es lista de canciones)
        categoria: Nombre de la categoría
        theme_colors: Dict con colores del tema
        output_file: Ruta del archivo de salida
    """
    
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 format
    prs.slide_height = Inches(7.5)
    
    # Determinar layout según número de canciones
    if not cards or len(cards) == 0:
        print('⚠️ No hay cartones para procesar')
        return None
    
    num_songs = len(cards[0])
    
    if num_songs <= 8:
        cartones_por_slide = 3
        layout_type = "horizontal_3"
    elif num_songs <= 12:
        cartones_por_slide = 2
        layout_type = "horizontal_2"
    else:
        cartones_por_slide = 1
        layout_type = "single"
    
    total_slides = (len(cards) + cartones_por_slide - 1) // cartones_por_slide
    
    print(f'📊 Creando {total_slides} diapositivas con {len(cards)} cartones ({num_songs} canciones/cartón, {cartones_por_slide} cartones/slide)...')
    
    for slide_idx in range(total_slides):
        # Crear slide en blanco
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Fondo de color suave
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors['background'])
        
        # Título de la slide
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        icon = theme_colors.get('icon', '🎵')
        title_frame.text = f"{icon} BINGO MUSICAL - {categoria}"
        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*theme_colors['title'])
        
        # Calcular cartones en esta slide
        start_idx = slide_idx * cartones_por_slide
        end_idx = min(start_idx + cartones_por_slide, len(cards))
        cards_in_slide = cards[start_idx:end_idx]
        
        # Configurar dimensiones según layout
        if layout_type == "horizontal_3":
            card_width = 3.0
            card_height = 5.5
            spacing = 0.2
            start_x = 0.5
            start_y = 1.2
            song_font_size = 11
            song_height = 0.7
        elif layout_type == "horizontal_2":
            card_width = 4.0
            card_height = 6.0
            spacing = 0.5
            start_x = 1.0
            start_y = 1.0
            song_font_size = 10
            song_height = 0.48
        else:  # single
            card_width = 6.0
            card_height = 6.5
            spacing = 0
            start_x = 2.0
            start_y = 0.8
            song_font_size = 9
            song_height = 0.28
        
        for card_pos, card in enumerate(cards_in_slide):
            card_number = start_idx + card_pos + 1
            x_pos = start_x + (card_width + spacing) * card_pos
            
            # Crear marco del cartón
            card_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x_pos), Inches(start_y),
                Inches(card_width), Inches(card_height)
            )
            
            # Estilo del marco
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card_shape.line.color.rgb = RGBColor(*theme_colors['border'])
            card_shape.line.width = Pt(3)
            
            # Título del cartón
            card_title = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(start_y + 0.1),
                Inches(card_width - 0.2), Inches(0.5)
            )
            card_title_frame = card_title.text_frame
            card_title_frame.text = f"Cartón {card_number}/{len(cards)}"
            card_title_p = card_title_frame.paragraphs[0]
            card_title_p.alignment = PP_ALIGN.CENTER
            card_title_p.font.size = Pt(14)
            card_title_p.font.bold = True
            card_title_p.font.color.rgb = RGBColor(*theme_colors['subtitle'])
            
            # Lista de canciones
            songs_y = start_y + 0.7
            
            for song_idx, song in enumerate(card):
                # Checkbox
                checkbox = slide.shapes.add_textbox(
                    Inches(x_pos + 0.15), Inches(songs_y + song_idx * song_height),
                    Inches(0.3), Inches(song_height - 0.05)
                )
                cb_frame = checkbox.text_frame
                cb_frame.text = "☐"
                cb_p = cb_frame.paragraphs[0]
                cb_p.font.size = Pt(18 if layout_type == "single" else 20)
                cb_p.font.color.rgb = RGBColor(*theme_colors['checkbox'])
                
                # Número y canción
                song_text = slide.shapes.add_textbox(
                    Inches(x_pos + 0.45), Inches(songs_y + song_idx * song_height),
                    Inches(card_width - 0.6), Inches(song_height - 0.05)
                )
                song_frame = song_text.text_frame
                song_frame.word_wrap = True
                
                # Limpiar título y artista
                song_clean = song.replace(' - Villancicos', '').replace(' - Canciones Infantiles', '')
                song_clean = song_clean.replace(' - Tradicional', '').replace(' - Pinkfong', '')
                song_clean = song_clean.replace(' - Canciones de la Granja', '').replace(' - Timbiriche', '')
                song_clean = song_clean.replace(' - Raphael', '').replace(' - Cri-Cri', '')
                
                song_frame.text = f"{song_idx + 1}. {song_clean}"
                song_p = song_frame.paragraphs[0]
                song_p.font.size = Pt(song_font_size)
                song_p.font.color.rgb = RGBColor(40, 40, 40)
                song_p.line_spacing = 0.9
        
        # Pie de página con número de slide
        footer = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.0),
            Inches(9), Inches(0.3)
        )
        footer_frame = footer.text_frame
        footer_frame.text = f"Diapositiva {slide_idx + 1} de {total_slides} · bingomusicalgratis.es"
        footer_p = footer_frame.paragraphs[0]
        footer_p.alignment = PP_ALIGN.CENTER
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(128, 128, 128)
        
        if (slide_idx + 1) % 5 == 0:
            print(f'  ✅ {slide_idx + 1}/{total_slides} diapositivas creadas...')
    
    # Guardar presentación
    prs.save(output_file)
    print(f'✅ Presentación guardada: {output_file}')
    
    return {
        'archivo': str(output_file),
        'slides': total_slides,
        'cartones': len(cards),
        'cartones_por_slide': cartones_por_slide,
        'layout': layout_type
    }

def load_cards_from_markdown(md_file):
    """Carga cartones desde archivo Markdown"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    cards = []
    current_card = []
    
    for line in content.split('\n'):
        if line.startswith('## Cartón'):
            if current_card:
                cards.append(current_card)
                current_card = []
        elif line and line[0].isdigit() and '. ' in line:
            # Extraer canción (quitar número)
            song = line.split('. ', 1)[1]
            current_card.append(song)
    
    if current_card:
        cards.append(current_card)
    
    return cards

def main():
    """Función principal"""
    print('\n🎨 Generador de PowerPoint - Cartones de Bingo Musical\n')
    
    base_path = Path(__file__).parent.parent / 'cartones'
    
    # Configuraciones de temas
    themes = {
        'navidad-pequeños': {
            'categoria': 'Navidad - Pequeños (8 canciones)',
            'icon': '🎄',
            'colors': {
                'background': (255, 250, 245),
                'title': (196, 30, 58),
                'subtitle': (34, 139, 34),
                'border': (196, 30, 58),
                'checkbox': (255, 215, 0)
            },
            'md_file': base_path / 'navidad' / 'pequeños' / 'cartones-navidad-pequeños.md',
            'output': base_path / 'navidad' / 'pequeños' / 'cartones-navidad-pequeños.pptx'
        },
        'navidad-medianos': {
            'categoria': 'Navidad - Medianos (12 canciones)',
            'icon': '🎄',
            'colors': {
                'background': (255, 250, 245),
                'title': (196, 30, 58),
                'subtitle': (34, 139, 34),
                'border': (196, 30, 58),
                'checkbox': (255, 215, 0)
            },
            'md_file': base_path / 'navidad' / 'medianos' / 'cartones-navidad-medianos.md',
            'output': base_path / 'navidad' / 'medianos' / 'cartones-navidad-medianos.pptx'
        },
        'navidad-grandes': {
            'categoria': 'Navidad - Grandes (20 canciones)',
            'icon': '🎄',
            'colors': {
                'background': (255, 250, 245),
                'title': (196, 30, 58),
                'subtitle': (34, 139, 34),
                'border': (196, 30, 58),
                'checkbox': (255, 215, 0)
            },
            'md_file': base_path / 'navidad' / 'grandes' / 'cartones-navidad-grandes.md',
            'output': base_path / 'navidad' / 'grandes' / 'cartones-navidad-grandes.pptx'
        },
        'clasicos-pop-pequeños': {
            'categoria': 'Clásicos Pop - Pequeños (8 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (255, 245, 250),
                'title': (255, 107, 157),
                'subtitle': (219, 39, 119),
                'border': (255, 107, 157),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'clasicos-del-pop' / 'pequeños' / 'cartones-clasicos-del-pop-pequeños.md',
            'output': base_path / 'clasicos-del-pop' / 'pequeños' / 'cartones-clasicos-del-pop-pequeños.pptx'
        },
        'clasicos-pop-medianos': {
            'categoria': 'Clásicos Pop - Medianos (12 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (255, 245, 250),
                'title': (255, 107, 157),
                'subtitle': (219, 39, 119),
                'border': (255, 107, 157),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'clasicos-del-pop' / 'medianos' / 'cartones-clasicos-del-pop-medianos.md',
            'output': base_path / 'clasicos-del-pop' / 'medianos' / 'cartones-clasicos-del-pop-medianos.pptx'
        },
        'clasicos-pop-grandes': {
            'categoria': 'Clásicos Pop - Grandes (20 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (255, 245, 250),
                'title': (255, 107, 157),
                'subtitle': (219, 39, 119),
                'border': (255, 107, 157),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'clasicos-del-pop' / 'grandes' / 'cartones-clasicos-del-pop-grandes.md',
            'output': base_path / 'clasicos-del-pop' / 'grandes' / 'cartones-clasicos-del-pop-grandes.pptx'
        },
        'pop-latino-pequeños': {
            'categoria': 'Pop Latino - Pequeños (8 canciones)',
            'icon': '💃',
            'colors': {
                'background': (255, 248, 240),
                'title': (255, 140, 66),
                'subtitle': (251, 146, 60),
                'border': (255, 140, 66),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'pop-latino-y-espanol' / 'pequeños' / 'cartones-pop-latino-y-espanol-pequeños.md',
            'output': base_path / 'pop-latino-y-espanol' / 'pequeños' / 'cartones-pop-latino-y-espanol-pequeños.pptx'
        },
        'pop-latino-medianos': {
            'categoria': 'Pop Latino - Medianos (12 canciones)',
            'icon': '💃',
            'colors': {
                'background': (255, 248, 240),
                'title': (255, 140, 66),
                'subtitle': (251, 146, 60),
                'border': (255, 140, 66),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'pop-latino-y-espanol' / 'medianos' / 'cartones-pop-latino-y-espanol-medianos.md',
            'output': base_path / 'pop-latino-y-espanol' / 'medianos' / 'cartones-pop-latino-y-espanol-medianos.pptx'
        },
        'pop-latino-grandes': {
            'categoria': 'Pop Latino - Grandes (20 canciones)',
            'icon': '💃',
            'colors': {
                'background': (255, 248, 240),
                'title': (255, 140, 66),
                'subtitle': (251, 146, 60),
                'border': (255, 140, 66),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'pop-latino-y-espanol' / 'grandes' / 'cartones-pop-latino-y-espanol-grandes.md',
            'output': base_path / 'pop-latino-y-espanol' / 'grandes' / 'cartones-pop-latino-y-espanol-grandes.pptx'
        },
        'cumpleanos-pequeños': {
            'categoria': 'Cumpleaños - Pequeños (8 canciones)',
            'icon': '🎂',
            'colors': {
                'background': (255, 254, 240),
                'title': (255, 217, 61),
                'subtitle': (251, 191, 36),
                'border': (255, 217, 61),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'cumpleanos' / 'pequeños' / 'cartones-cumpleanos-pequeños.md',
            'output': base_path / 'cumpleanos' / 'pequeños' / 'cartones-cumpleanos-pequeños.pptx'
        },
        'cumpleanos-medianos': {
            'categoria': 'Cumpleaños - Medianos (12 canciones)',
            'icon': '🎂',
            'colors': {
                'background': (255, 254, 240),
                'title': (255, 217, 61),
                'subtitle': (251, 191, 36),
                'border': (255, 217, 61),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'cumpleanos' / 'medianos' / 'cartones-cumpleanos-medianos.md',
            'output': base_path / 'cumpleanos' / 'medianos' / 'cartones-cumpleanos-medianos.pptx'
        },
        'otono-pequeños': {
            'categoria': 'Otoño - Pequeños (8 canciones)',
            'icon': '🍂',
            'colors': {
                'background': (255, 247, 237),
                'title': (212, 165, 116),
                'subtitle': (180, 130, 70),
                'border': (212, 165, 116),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-de-otono' / 'pequeños' / 'cartones-musica-de-otono-pequeños.md',
            'output': base_path / 'musica-de-otono' / 'pequeños' / 'cartones-musica-de-otono-pequeños.pptx'
        },
        'otono-medianos': {
            'categoria': 'Otoño - Medianos (12 canciones)',
            'icon': '🍂',
            'colors': {
                'background': (255, 247, 237),
                'title': (212, 165, 116),
                'subtitle': (180, 130, 70),
                'border': (212, 165, 116),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-de-otono' / 'medianos' / 'cartones-musica-de-otono-medianos.md',
            'output': base_path / 'musica-de-otono' / 'medianos' / 'cartones-musica-de-otono-medianos.pptx'
        },
        'villancicos-infantil': {
            'categoria': 'Villancicos Infantil',
            'icon': '🎄',
            'colors': {
                'background': (255, 250, 245),
                'title': (196, 30, 58),
                'subtitle': (34, 139, 34),
                'border': (196, 30, 58),
                'checkbox': (255, 215, 0)
            },
            'md_file': base_path / 'villancicos-infantil' / 'cartones-villancicos-infantil.md',
            'output': base_path / 'villancicos-infantil' / 'cartones-villancicos-infantil.pptx'
        },
        'infantil': {
            'categoria': 'Canciones Infantiles',
            'icon': '🎈',
            'colors': {
                'background': (245, 250, 255),
                'title': (99, 102, 241),
                'subtitle': (168, 85, 247),
                'border': (99, 102, 241),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'infantil' / 'cartones-infantil.md',
            'output': base_path / 'infantil' / 'cartones-infantil.pptx'
        },
        'rock-pequeños': {
            'categoria': 'Rock - Pequeños (8 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock' / 'pequeños' / 'cartones-rock-pequeños.md',
            'output': base_path / 'rock' / 'pequeños' / 'cartones-rock-pequeños.pptx'
        },
        'rock-medianos': {
            'categoria': 'Rock - Medianos (12 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock' / 'medianos' / 'cartones-rock-medianos.md',
            'output': base_path / 'rock' / 'medianos' / 'cartones-rock-medianos.pptx'
        },
        'rock-grandes': {
            'categoria': 'Rock - Grandes (20 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock' / 'grandes' / 'cartones-rock-grandes.md',
            'output': base_path / 'rock' / 'grandes' / 'cartones-rock-grandes.pptx'
        },
        'rock-clasico-pequeños': {
            'categoria': 'Rock Clásico - Pequeños (8 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock-clasico' / 'pequeños' / 'cartones-rock-clasico-pequeños.md',
            'output': base_path / 'rock-clasico' / 'pequeños' / 'cartones-rock-clasico-pequeños.pptx'
        },
        'rock-clasico-medianos': {
            'categoria': 'Rock Clásico - Medianos (12 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock-clasico' / 'medianos' / 'cartones-rock-clasico-medianos.md',
            'output': base_path / 'rock-clasico' / 'medianos' / 'cartones-rock-clasico-medianos.pptx'
        },
        'rock-clasico-grandes': {
            'categoria': 'Rock Clásico - Grandes (20 canciones)',
            'icon': '🎸',
            'colors': {
                'background': (250, 245, 250),
                'title': (138, 43, 226),
                'subtitle': (148, 0, 211),
                'border': (138, 43, 226),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'rock-clasico' / 'grandes' / 'cartones-rock-clasico-grandes.md',
            'output': base_path / 'rock-clasico' / 'grandes' / 'cartones-rock-clasico-grandes.pptx'
        },
        'espanol-pequeños': {
            'categoria': 'Música en Español - Pequeños (8 canciones)',
            'icon': '🇪🇸',
            'colors': {
                'background': (255, 250, 240),
                'title': (220, 38, 38),
                'subtitle': (239, 68, 68),
                'border': (220, 38, 38),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-espanol' / 'pequeños' / 'cartones-musica-en-espanol-pequeños.md',
            'output': base_path / 'musica-en-espanol' / 'pequeños' / 'cartones-musica-en-espanol-pequeños.pptx'
        },
        'espanol-medianos': {
            'categoria': 'Música en Español - Medianos (12 canciones)',
            'icon': '🇪🇸',
            'colors': {
                'background': (255, 250, 240),
                'title': (220, 38, 38),
                'subtitle': (239, 68, 68),
                'border': (220, 38, 38),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-espanol' / 'medianos' / 'cartones-musica-en-espanol-medianos.md',
            'output': base_path / 'musica-en-espanol' / 'medianos' / 'cartones-musica-en-espanol-medianos.pptx'
        },
        'espanol-grandes': {
            'categoria': 'Música en Español - Grandes (20 canciones)',
            'icon': '🇪🇸',
            'colors': {
                'background': (255, 250, 240),
                'title': (220, 38, 38),
                'subtitle': (239, 68, 68),
                'border': (220, 38, 38),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-espanol' / 'grandes' / 'cartones-musica-en-espanol-grandes.md',
            'output': base_path / 'musica-en-espanol' / 'grandes' / 'cartones-musica-en-espanol-grandes.pptx'
        },
        'ingles-pequeños': {
            'categoria': 'Música en Inglés - Pequeños (8 canciones)',
            'icon': '🇬🇧',
            'colors': {
                'background': (240, 248, 255),
                'title': (30, 64, 175),
                'subtitle': (59, 130, 246),
                'border': (30, 64, 175),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-ingles' / 'pequeños' / 'cartones-musica-en-ingles-pequeños.md',
            'output': base_path / 'musica-en-ingles' / 'pequeños' / 'cartones-musica-en-ingles-pequeños.pptx'
        },
        'ingles-medianos': {
            'categoria': 'Música en Inglés - Medianos (12 canciones)',
            'icon': '🇬🇧',
            'colors': {
                'background': (240, 248, 255),
                'title': (30, 64, 175),
                'subtitle': (59, 130, 246),
                'border': (30, 64, 175),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-ingles' / 'medianos' / 'cartones-musica-en-ingles-medianos.md',
            'output': base_path / 'musica-en-ingles' / 'medianos' / 'cartones-musica-en-ingles-medianos.pptx'
        },
        'ingles-grandes': {
            'categoria': 'Música en Inglés - Grandes (20 canciones)',
            'icon': '🇬🇧',
            'colors': {
                'background': (240, 248, 255),
                'title': (30, 64, 175),
                'subtitle': (59, 130, 246),
                'border': (30, 64, 175),
                'checkbox': (252, 211, 77)
            },
            'md_file': base_path / 'musica-en-ingles' / 'grandes' / 'cartones-musica-en-ingles-grandes.md',
            'output': base_path / 'musica-en-ingles' / 'grandes' / 'cartones-musica-en-ingles-grandes.pptx'
        }
    }
    
    # Generar todas las presentaciones
    results = []
    
    for theme_key, theme_config in themes.items():
        print(f'\n📁 Procesando: {theme_config["categoria"]}')
        
        # Cargar cartones
        print(f'   Cargando cartones desde MD...')
        cards = load_cards_from_markdown(theme_config['md_file'])
        print(f'   ✅ {len(cards)} cartones cargados')
        
        # Crear PowerPoint
        result = create_bingo_pptx(
            cards=cards,
            categoria=theme_config['categoria'],
            theme_colors=theme_config['colors'],
            output_file=theme_config['output']
        )
        results.append(result)
    
    # Resumen final
    print(f'\n🎉 ¡Generación completada!\n')
    print('📊 Resumen:')
    for idx, result in enumerate(results, 1):
        print(f'   {idx}. {Path(result["archivo"]).name}')
        print(f'      - {result["slides"]} diapositivas')
        print(f'      - {result["cartones"]} cartones')
        print(f'      - {result["cartones_por_slide"]} cartones por slide')
    print()

if __name__ == '__main__':
    main()
