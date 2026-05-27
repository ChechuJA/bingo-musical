"""Utilidades para generar PPTX de cartones.

Este módulo existe para reutilizar el generador PPTX desde distintos scripts
(sin duplicar lógica).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_bingo_pptx(
    cards: list[list[str]],
    categoria: str,
    theme_colors: dict,
    output_file: Path | str,
    *,
    card_style: str = "list",
    grid_wildcard_text: str = "⭐ COMODÍN",
    grid_wildcard_image: Path | str | None = None,
    force_cards_per_slide: int | None = None,
    logo_path: Path | str | None = None,
    official_logo_path: Path | str | None = None,
    logo_width: float = 0.8,
    official_logo_scale: float = 0.8,
    custom_logo_scale: float = 1.5,
    logo_position: str = "top-left",
    show_slide_title: bool = True,
    show_card_number: bool = True,
    logo_on_each_card: bool = False,
    footer_text: str = "bingomusicalgratis.es",
    font_title: str | None = None,
    font_artist: str | None = None,
):
    """Crea presentación PowerPoint con cartones de bingo.

    Layout dinámico según número de canciones:
    - 6-8 canciones: 3 cartones por slide
    - 10-12 canciones: 2 cartones por slide
    - 15-20 canciones: 1 cartón por slide

    Estilos de cartón:
    - card_style="list": lista vertical con checkboxes (estilo actual)
    - card_style="grid3x3": cuadrícula 3x3 (pensado para 8 canciones + 1 comodín)

    Logo:
    - official_logo_path: logo oficial (izquierda)
    - logo_path: logo personalizado (derecha)
    - logo_width: ancho/alto base para logos en pulgadas (default 0.8)
    - custom_logo_scale: multiplicador del logo personalizado (default 1.5)
    - official_logo_scale: multiplicador del logo oficial (default 0.8)
    - logo_position: "top-left", "top-right", "top-center" (default "top-left")
    """

    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9
    prs.slide_height = Inches(7.5)

    if not cards:
        print("⚠️ No hay cartones para procesar")
        return None

    num_songs = len(cards[0])

    if force_cards_per_slide is not None:
        if force_cards_per_slide not in (1, 2, 3, 4):
            raise ValueError("force_cards_per_slide debe ser 1, 2, 3 o 4")
        cartones_por_slide = force_cards_per_slide
        layout_type = {4: "grid_2x2", 3: "horizontal_3", 2: "horizontal_2", 1: "single"}[force_cards_per_slide]
    else:
        if num_songs <= 8:
            cartones_por_slide = 3
            layout_type = "horizontal_3"
        elif num_songs <= 12:
            cartones_por_slide = 2
            layout_type = "horizontal_2"
        else:
            cartones_por_slide = 1
            layout_type = "single"

    total_slides = (len(cards) + cartones_por_slide - 1) // cartones_por_slide

    print(
        f"📊 Creando {total_slides} diapositivas con {len(cards)} cartones "
        f"({num_songs} canciones/cartón, {cartones_por_slide} cartones/slide)..."
    )

    for slide_idx in range(total_slides):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["background"])

        # --- Logo a nivel de slide (solo si NO va por cartón) ---
        if not logo_on_each_card:
            _official_logo = Path(official_logo_path) if official_logo_path else None
            _custom_logo = Path(logo_path) if logo_path else None

            custom_logo_height = logo_width * custom_logo_scale
            official_logo_height = logo_width * official_logo_scale
            if _official_logo and _official_logo.exists() and _custom_logo and _custom_logo.exists():
                # Izquierda: logo personalizado (cliente)
                slide.shapes.add_picture(
                    str(_custom_logo),
                    Inches(0.2),
                    Inches(0.1),
                    height=Inches(custom_logo_height),
                )
                # Derecha: logo oficial
                pic_official = slide.shapes.add_picture(
                    str(_official_logo),
                    Inches(0),
                    Inches(0.1),
                    height=Inches(official_logo_height),
                )
                pic_official.left = int(Inches(10 - 0.2) - pic_official.width)
            elif _custom_logo and _custom_logo.exists():
                custom_logo_size = logo_width * custom_logo_scale
                if logo_position == "top-right":
                    lx = Inches(10 - custom_logo_size - 0.2)
                elif logo_position == "top-center":
                    lx = Inches((10 - custom_logo_size) / 2)
                else:  # top-left
                    lx = Inches(0.2)
                slide.shapes.add_picture(
                    str(_custom_logo), lx, Inches(0.1), height=Inches(custom_logo_size)
                )
            elif _official_logo and _official_logo.exists():
                slide.shapes.add_picture(
                    str(_official_logo),
                    Inches(0.2),
                    Inches(0.1),
                    height=Inches(official_logo_height),
                )

        if show_slide_title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            icon = theme_colors.get("icon", "🎵")
            title_frame.text = f"{icon} BINGO MUSICAL - {categoria}"
            title_p = title_frame.paragraphs[0]
            title_p.alignment = PP_ALIGN.CENTER
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            title_p.font.color.rgb = RGBColor(*theme_colors["title"])

        start_idx = slide_idx * cartones_por_slide
        end_idx = min(start_idx + cartones_por_slide, len(cards))
        cards_in_slide = cards[start_idx:end_idx]

        if layout_type == "grid_2x2":
            card_width = 4.75
            card_height = 3.0
            spacing_x = 0.18
            spacing_y = 0.1
            spacing = 0.18
            start_x = 0.15
            start_y = 0.05
            song_font_size = 10
            song_height = 0.4
        elif layout_type == "horizontal_3":
            card_width = 3.0
            card_height = 5.5
            spacing = 0.2
            start_x = 0.5
            start_y = 1.2 if show_slide_title else 0.15
            song_font_size = 10
            song_height = 0.62
        elif layout_type == "horizontal_2":
            card_width = 4.7
            card_height = 7.2
            spacing = 0.25
            start_x = 0.15
            start_y = 0.15 if not show_slide_title else 1.0
            song_font_size = 10
            song_height = 0.48
        else:
            card_width = 6.0
            card_height = 6.5
            spacing = 0
            start_x = 2.0
            start_y = 0.8 if show_slide_title else 0.15
            song_font_size = 9
            song_height = 0.28

        for card_pos, card in enumerate(cards_in_slide):
            card_number = start_idx + card_pos + 1

            if layout_type == "grid_2x2":
                col = card_pos % 2
                row = card_pos // 2
                x_pos = start_x + col * (card_width + spacing_x)
                y_pos = start_y + row * (card_height + spacing_y)
            else:
                x_pos = start_x + (card_width + spacing) * card_pos
                y_pos = start_y

            card_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x_pos),
                Inches(y_pos),
                Inches(card_width),
                Inches(card_height),
            )

            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card_shape.line.color.rgb = RGBColor(*theme_colors["border"])
            card_shape.line.width = Pt(3)

            # --- Logo por cartón ---
            logo_area_h = 0.0
            if logo_on_each_card:
                _official_logo = Path(official_logo_path) if official_logo_path else None
                _custom_logo = Path(logo_path) if logo_path else None
                has_official = bool(_official_logo and _official_logo.exists())
                has_custom = bool(_custom_logo and _custom_logo.exists())
                custom_logo_height = logo_width * custom_logo_scale
                official_logo_height = logo_width * official_logo_scale

                if has_official and has_custom:
                    logo_area_h = max(official_logo_height, custom_logo_height) + 0.12
                    left_x = x_pos + 0.12
                    right_anchor = Inches(x_pos + card_width - 0.12)
                    logo_y = y_pos + 0.08
                    # Izquierda: logo personalizado (cliente)
                    slide.shapes.add_picture(
                        str(_custom_logo), Inches(left_x), Inches(logo_y), height=Inches(custom_logo_height)
                    )
                    # Derecha: logo oficial
                    pic_official = slide.shapes.add_picture(
                        str(_official_logo), Inches(0), Inches(logo_y), height=Inches(official_logo_height)
                    )
                    pic_official.left = int(right_anchor - pic_official.width)
                elif has_custom:
                    logo_area_h = custom_logo_height + 0.1
                    _lh = custom_logo_height
                    _lx = x_pos + (card_width - _lh) / 2
                    _ly = y_pos + 0.1
                    slide.shapes.add_picture(
                        str(_custom_logo), Inches(_lx), Inches(_ly), height=Inches(_lh)
                    )
                elif has_official:
                    logo_area_h = official_logo_height + 0.1
                    _lh = official_logo_height
                    _lx = x_pos + 0.12
                    _ly = y_pos + 0.1
                    slide.shapes.add_picture(
                        str(_official_logo), Inches(_lx), Inches(_ly), height=Inches(_lh)
                    )

            if show_card_number:
                card_title = slide.shapes.add_textbox(
                    Inches(x_pos + 0.1),
                    Inches(y_pos + logo_area_h + 0.1),
                    Inches(card_width - 0.2),
                    Inches(0.45),
                )
                card_title_frame = card_title.text_frame
                card_title_frame.text = f"Cartón {card_number}/{len(cards)}"
                card_title_p = card_title_frame.paragraphs[0]
                card_title_p.alignment = PP_ALIGN.CENTER
                card_title_p.font.size = Pt(14)
                card_title_p.font.bold = True
                card_title_p.font.color.rgb = RGBColor(*theme_colors["subtitle"])
                header_h = logo_area_h + 0.6
            else:
                header_h = logo_area_h + 0.05

            def _clean_song(s: str) -> str:
                return (
                    s.replace(" - Villancicos", "")
                    .replace(" - Canciones Infantiles", "")
                    .replace(" - Tradicional", "")
                    .replace(" - Pinkfong", "")
                    .replace(" - Canciones de la Granja", "")
                    .replace(" - Timbiriche", "")
                    .replace(" - Raphael", "")
                    .replace(" - Cri-Cri", "")
                )

            if card_style == "grid" and card_style != "grid3x3":
                # Cuadrícula genérica: 3 columnas, filas calculadas automáticamente
                n_songs = len(card)
                cols = 3
                rows = (n_songs + cols - 1) // cols

                footer_h = 0.20
                card_num_h = 0.18
                margin = 0.05
                grid_x = x_pos + margin
                grid_y = y_pos + header_h
                grid_w = card_width - (margin * 2)
                grid_h = card_height - header_h - footer_h - card_num_h - 0.04
                cell_w = grid_w / cols
                cell_h = grid_h / rows

                # Fuente adaptativa al tamaño de celda
                cell_font_title = min(10, max(7, int(cell_h * 9)))
                cell_font_artist = max(6, cell_font_title - 2)

                def _split_song(s: str) -> tuple[str, str]:
                    """Separa 'Título - Artista' en dos cadenas."""
                    clean = _clean_song(s)
                    if " - " in clean:
                        parts = clean.split(" - ", 1)
                        return parts[0].strip(), parts[1].strip()
                    return clean, ""

                for idx, song in enumerate(card):
                    r = idx // cols
                    c = idx % cols
                    cx = grid_x + c * cell_w
                    cy = grid_y + r * cell_h

                    cell_shape = slide.shapes.add_shape(
                        1,
                        Inches(cx),
                        Inches(cy),
                        Inches(cell_w),
                        Inches(cell_h),
                    )
                    cell_shape.fill.solid()
                    cell_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    cell_shape.line.color.rgb = RGBColor(*theme_colors["border"])
                    cell_shape.line.width = Pt(1.0)

                    title_str, artist_str = _split_song(song)

                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.03),
                        Inches(cy + 0.03),
                        Inches(cell_w - 0.06),
                        Inches(cell_h - 0.06),
                    )
                    tf = tb.text_frame
                    tf.word_wrap = True

                    # Línea 1: Título en negrita
                    p1 = tf.paragraphs[0]
                    p1.alignment = PP_ALIGN.CENTER
                    run1 = p1.add_run()
                    run1.text = title_str
                    run1.font.bold = True
                    run1.font.size = Pt(cell_font_title)
                    run1.font.color.rgb = RGBColor(20, 20, 20)
                    if font_title:
                        run1.font.name = font_title

                    # Línea 2: Artista en gris más pequeño
                    if artist_str:
                        p2 = tf.add_paragraph()
                        p2.alignment = PP_ALIGN.CENTER
                        run2 = p2.add_run()
                        run2.text = artist_str
                        run2.font.bold = False
                        run2.font.size = Pt(cell_font_artist)
                        run2.font.color.rgb = RGBColor(110, 110, 110)
                        if font_artist:
                            run2.font.name = font_artist

                # Número de cartón abajo a la derecha: Cartón #034
                cn_box = slide.shapes.add_textbox(
                    Inches(x_pos + card_width - 1.15),
                    Inches(y_pos + card_height - card_num_h - 0.01),
                    Inches(1.12),
                    Inches(card_num_h),
                )
                cn_f = cn_box.text_frame
                cn_f.text = f"Cartón #{card_number:03d}"
                cn_p = cn_f.paragraphs[0]
                cn_p.alignment = PP_ALIGN.RIGHT
                cn_p.font.size = Pt(7)
                cn_p.font.color.rgb = RGBColor(160, 160, 160)

                # Footer web izquierda (no choca con numeración)
                if footer_text:
                    ft = slide.shapes.add_textbox(
                        Inches(x_pos + 0.05),
                        Inches(y_pos + card_height - footer_h - 0.01),
                        Inches(card_width - 1.2),
                        Inches(footer_h),
                    )
                    ft_f = ft.text_frame
                    ft_f.text = footer_text
                    ft_p = ft_f.paragraphs[0]
                    ft_p.alignment = PP_ALIGN.LEFT
                    ft_p.font.size = Pt(7)
                    ft_p.font.color.rgb = RGBColor(180, 180, 180)

            elif card_style == "grid3x3":
                # Cuadrícula 3x3: 8 canciones + 1 comodín
                if len(card) > 8:
                    raise ValueError(
                        f"grid3x3 solo soporta hasta 8 canciones por cartón. Recibido: {len(card)}"
                    )

                margin = 0.12
                grid_x = x_pos + margin
                grid_y = start_y + 0.68
                grid_w = card_width - (margin * 2)
                grid_h = card_height - 0.9
                cols = 3
                rows = 3
                cell_w = grid_w / cols
                cell_h = grid_h / rows

                # Contenido de 9 celdas (última = comodín)
                cell_items: list[dict] = []
                for song in card:
                    cell_items.append({"type": "song", "value": _clean_song(song)})
                while len(cell_items) < 8:
                    cell_items.append({"type": "empty", "value": ""})
                cell_items.append({"type": "wildcard", "value": grid_wildcard_text})

                for idx, item in enumerate(cell_items):
                    r = idx // cols
                    c = idx % cols
                    cx = grid_x + c * cell_w
                    cy = grid_y + r * cell_h

                    cell_shape = slide.shapes.add_shape(
                        1,
                        Inches(cx),
                        Inches(cy),
                        Inches(cell_w),
                        Inches(cell_h),
                    )
                    cell_shape.fill.solid()
                    cell_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    cell_shape.line.color.rgb = RGBColor(*theme_colors["border"])
                    cell_shape.line.width = Pt(1.5)

                    # Wildcard: imagen (si existe) o texto
                    if item["type"] == "wildcard" and grid_wildcard_image:
                        img_path = Path(grid_wildcard_image)
                        if img_path.exists():
                            # Ajuste simple: centrar dentro de la celda
                            pad = 0.08
                            slide.shapes.add_picture(
                                str(img_path),
                                Inches(cx + pad),
                                Inches(cy + pad),
                                width=Inches(cell_w - (pad * 2)),
                                height=Inches(cell_h - (pad * 2)),
                            )
                            continue

                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.06),
                        Inches(cy + 0.06),
                        Inches(cell_w - 0.12),
                        Inches(cell_h - 0.12),
                    )
                    tf = tb.text_frame
                    tf.word_wrap = True
                    tf.text = item["value"]
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = item["type"] == "wildcard"
                    if layout_type == "horizontal_3":
                        grid_font = 10
                    elif layout_type == "horizontal_2":
                        grid_font = 12
                    else:
                        grid_font = 14
                    p.font.size = Pt(grid_font)
                    p.font.color.rgb = RGBColor(40, 40, 40)
            else:
                # Lista vertical con checkbox
                songs_y = y_pos + header_h

                for song_idx, song in enumerate(card):
                    checkbox = slide.shapes.add_textbox(
                        Inches(x_pos + 0.15),
                        Inches(songs_y + song_idx * song_height),
                        Inches(0.3),
                        Inches(song_height - 0.05),
                    )
                    cb_frame = checkbox.text_frame
                    cb_frame.text = "☐"
                    cb_p = cb_frame.paragraphs[0]
                    cb_p.font.size = Pt(18 if layout_type == "single" else 18)
                    cb_p.font.color.rgb = RGBColor(*theme_colors["checkbox"])

                    song_text = slide.shapes.add_textbox(
                        Inches(x_pos + 0.45),
                        Inches(songs_y + song_idx * song_height),
                        Inches(card_width - 0.6),
                        Inches(song_height - 0.05),
                    )
                    song_frame = song_text.text_frame
                    song_frame.word_wrap = True

                    song_frame.text = f"{song_idx + 1}. {_clean_song(song)}"
                    song_p = song_frame.paragraphs[0]
                    song_p.font.size = Pt(song_font_size)
                    song_p.font.color.rgb = RGBColor(40, 40, 40)
                    song_p.line_spacing = 0.9

        if not (card_style == "grid" and not show_card_number):
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(9), Inches(0.3))
            footer_frame = footer.text_frame
            footer_frame.text = f"Diapositiva {slide_idx + 1} de {total_slides} · bingomusicalgratis.es"
            footer_p = footer_frame.paragraphs[0]
            footer_p.alignment = PP_ALIGN.CENTER
            footer_p.font.size = Pt(10)
            footer_p.font.color.rgb = RGBColor(128, 128, 128)

        if (slide_idx + 1) % 5 == 0:
            print(f"  ✅ {slide_idx + 1}/{total_slides} diapositivas creadas...")

    prs.save(output_file)
    print(f"✅ Presentación guardada: {output_file}")

    return {
        "archivo": str(output_file),
        "slides": total_slides,
        "cartones": len(cards),
        "cartones_por_slide": cartones_por_slide,
        "layout": layout_type,
    }


def load_cards_from_markdown(md_file: Path | str) -> list[list[str]]:
    """Carga cartones desde archivo Markdown."""
    md_file = Path(md_file)
    content = md_file.read_text(encoding="utf-8")

    cards: list[list[str]] = []
    current_card: list[str] = []

    for line in content.split("\n"):
        if line.startswith("## Cartón"):
            if current_card:
                cards.append(current_card)
                current_card = []
        elif line and line[0].isdigit() and ". " in line:
            song = line.split(". ", 1)[1]
            current_card.append(song)

    if current_card:
        cards.append(current_card)

    return cards
